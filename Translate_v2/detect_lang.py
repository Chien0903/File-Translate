import os
import docx
import openpyxl
import pypdf as PyPDF2
import pptx
import zipfile
import xml.etree.ElementTree as ET
from fastapi import HTTPException
from google.cloud import translate
from dotenv import load_dotenv
try:
    from .create_glossary import language_pair
except ImportError:
    from create_glossary import language_pair
try:
    import mammoth  # type: ignore
except Exception:
    mammoth = None

load_dotenv()
# Định nghĩa các ngôn ngữ hỗ trợ
VIETNAMESE_SYMBOL = 'vi'
JAPANESE_SYMBOL = 'ja'
CHINESE_SIMPLIFIED_SYMBOL = 'zh-CN'

CHINESE_TRADITIONAL_SYMBOL = 'zh-TW'
ENGLISH_SYMBOL = 'en'

# Lấy giá trị biến môi trường
google_credentials_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
PROJECT_ID = os.getenv("PROJECT_ID")
if google_credentials_path:
    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = google_credentials_path

# Định nghĩa các ngôn ngữ hỗ trợ (mở rộng)
LANGUAGES = {
    'vi': "Vietnamese",
    'ja': "Japanese",
    'zh-CN': "Chinese (Simplified)",
    'zh-TW': "Chinese (Traditional)",
    'en': "English",
    'th': "Thai",
    'bn': "Bengali",
    'hi': "Hindi",
    'id': "Indonesian",
    'or': "Oriya",
}

# language_pair được import từ create_glossary.py để đồng bộ với backend
def extract_content(file_path: str) -> str:
    """Trích xuất nội dung từ file (PDF, DOCX, XLSX) để phát hiện ngôn ngữ"""
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    file_extension = file_path.split(".")[-1].lower()

    if file_extension == "pdf":
        return extract_from_pdf(file_path)
    elif file_extension == "docx":
        return extract_from_docx(file_path)
    elif file_extension == "xlsx":
        return extract_from_xlsx(file_path)
    elif file_extension == "pptx":
        return extract_from_pptx(file_path)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")
def extract_from_pptx(file_path: str) -> str:
    """Trích xuất văn bản từ file PPTX"""
    content = ""
    prs = pptx.Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + " "
                if len(content.split()) >= 1000:
                    break
        if len(content.split()) >= 1000:
            break
    return get_first_1000_words(content)
def get_first_1000_words(text: str) -> str:
    """Lấy 1000 từ đầu tiên của văn bản (bỏ qua các từ chứa số)"""
    words = text.split()
    non_numeric_words = [w for w in words if not any(c.isdigit() for c in w)]
    return ' '.join(non_numeric_words[:1000])

def extract_from_pdf(file_path: str) -> str:
    """Trích xuất văn bản từ file PDF"""
    content = ""
    with open(file_path, "rb") as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page in pdf_reader.pages:
            content += page.extract_text() or ""
            if len(content.split()) >= 1000:
                break
    return get_first_1000_words(content)

def extract_from_docx(file_path: str) -> str:
    """
    Trích xuất văn bản từ file DOCX.
    Ưu tiên python-docx (bao gồm cả header/footer), fallback sang mammoth nếu lỗi.
    Giới hạn tối đa 1000 từ.
    """
    content = ""
    try:
        # Thử đọc bằng python-docx
        doc = docx.Document(file_path)

        # Lấy text từ body qua OOXML <w:t>
        for node in doc.element.body.iter():
            if node.tag.endswith('}t') and node.text:
                content += node.text + " "
                if len(content.split()) >= 1000:
                    return get_first_1000_words(content)

        # Lấy text từ header/footer nếu có
        try:
            if hasattr(doc.part, "header_parts"):
                for hp in doc.part.header_parts:
                    for node in hp._element.iter():
                        if node.tag.endswith('}t') and getattr(node, "text", None):
                            content += node.text + " "
                            if len(content.split()) >= 1000:
                                return get_first_1000_words(content)
            if hasattr(doc.part, "footer_parts"):
                for fp in doc.part.footer_parts:
                    for node in fp._element.iter():
                        if node.tag.endswith('}t') and getattr(node, "text", None):
                            content += node.text + " "
                            if len(content.split()) >= 1000:
                                return get_first_1000_words(content)
        except Exception:
            pass

        return get_first_1000_words(content)

    except Exception as e:
        # Fallback sang mammoth nếu python-docx lỗi
        try:
            if mammoth is None:
                return ""
            with open(file_path, "rb") as f:
                result = mammoth.extract_raw_text(f)
                text = result.value or ""
            return get_first_1000_words(text)
        except Exception:
            return ""

def extract_from_xlsx(file_path: str) -> str:
    """
    Trích xuất văn bản từ file XLSX.
    - Duyệt nhiều sheet, lấy giá trị hiển thị (data_only=True).
    - Bổ sung quét text trong shapes/textboxes (xl/drawings/*.xml) nếu còn thiếu.
    """
    content = ""
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    content += str(cell) + " "
                    if len(content.split()) >= 1000:
                        return get_first_1000_words(content)
            if len(content.split()) >= 1000:
                return get_first_1000_words(content)

    # Nếu chưa đủ, đọc như zip để lấy text trong drawings (textbox/shape)
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            for name in z.namelist():
                if name.startswith("xl/drawings/") and name.endswith(".xml"):
                    data = z.read(name)
                    try:
                        root = ET.fromstring(data)
                        for el in root.iter():
                            if el.tag.endswith('}t') and el.text:
                                content += el.text + " "
                                if len(content.split()) >= 1000:
                                    return get_first_1000_words(content)
                    except ET.ParseError:
                        continue
    except Exception:
        pass

    return get_first_1000_words(content)

def detect_language(text: str):
    """Phát hiện ngôn ngữ của file"""
    try:
        client = translate.TranslationServiceClient()
        parent = f"projects/{PROJECT_ID}/locations/us-central1"
        request = translate.DetectLanguageRequest(content=text, parent=parent)
        response = client.detect_language(request=request)
        detected_language = response.languages[0].language_code
        return detected_language
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error detecting language: {str(e)}")
